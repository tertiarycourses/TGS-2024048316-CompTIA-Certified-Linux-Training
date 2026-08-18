#!/usr/bin/env python3
"""build_slides.py — WSQ course slide deck for TGS-2024048316 (CompTIA Linux+ XK0-006).

All-white Tertiary house style, built from the mandatory wsq-slides **visual component
library** (tile grids, horizontal flow diagrams, colour-topped cards, profile cards,
activity/step/verify slides, section dividers) — never bullet walls. Content is driven
entirely by the single source course_content.py so the deck stays 100% aligned with the
Lesson Plan, Learner Guide (DOCX + MD) and the 30 labs.

Also writes courseware/slide_map.json (lab/domain -> deck page) so the Lesson Plan can
cite exact slide numbers.

Component helpers ported from the wsq-slides reference pipeline.
"""
import os, sys, json
from pptx import Presentation
from pptx.util import Inches as _EmuInches, Pt as _EmuPt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

_HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.join(_HERE, os.pardir, "wsq-learner-guide"))
from course_content import (COURSE, DOMAINS, LEARNING_OUTCOMES, DOMAIN_CONCEPTS,
                            DOMAIN_OUTCOMES, domain_labs, COURSEWARE, ASSETS,
                            REFERENCE_DECK, REFERENCE_CHAPTERS, DOMAIN_CHAPTERS,
                            REFERENCE_DROP)

# The house components are authored on a 13.333 x 7.5 in canvas; the merged deck
# keeps the legacy deck's 10 x 5.625 in canvas (same 16:9), so every house
# coordinate and font size is scaled uniformly by 0.75.
SCALE = 0.75
def Inches(v): return _EmuInches(v * SCALE)
def Pt(v): return _EmuPt(v * SCALE)

# ---------------- palette ----------------
BLUE=RGBColor(0x1F,0x6F,0xEB); TEAL=RGBColor(0x10,0xB9,0x81); AMBER=RGBColor(0xF5,0x9E,0x0B)
INK=RGBColor(0x16,0x1B,0x26); GREY=RGBColor(0x5B,0x63,0x72); LIGHT=RGBColor(0xF5,0xF8,0xFC)
WHITE=RGBColor(0xFF,0xFF,0xFF); LINE=RGBColor(0xE2,0xE8,0xF0); VIOLET=RGBColor(0x7C,0x3A,0xED)
RED=RGBColor(0xC8,0x10,0x2E)
PALETTE=[BLUE,TEAL,VIOLET,AMBER]

# Open the LEGACY v3 deck as the base package: all 22 chapters of teaching
# slides (text + graphics) stay in the file untouched; the house slides are
# appended and everything is re-sequenced into exam-domain order at the end.
prs=Presentation(REFERENCE_DECK)
SW,SH=prs.slide_width,prs.slide_height
BLANK=None
for _m in prs.slide_masters:
    for _l in _m.slide_layouts:
        if _l.name.strip().upper()=="BLANK": BLANK=_l; break
    if BLANK is not None: break
assert BLANK is not None, "no BLANK layout found in the reference deck"

# The legacy masters/layouts carry a footer with a WRONG UEN ("20120096W");
# correct it everywhere, and strip the legacy footer entirely from the BLANK
# layout so the house slides (which draw their own footer) don't inherit it.
_BAD_UEN, _GOOD_UEN = "20120096W", "201200696W"
def _fix_footer(container, strip=False):
    for sh in list(container.shapes):
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text
        if _BAD_UEN in t or "This material belongs to" in t:
            if strip:
                sh._element.getparent().remove(sh._element)
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if _BAD_UEN in r.text:
                        r.text = r.text.replace(_BAD_UEN, _GOOD_UEN)
for _m in prs.slide_masters:
    _fix_footer(_m)
    for _l in _m.slide_layouts:
        _fix_footer(_l, strip=(_l is BLANK))
for _s in prs.slides:
    _fix_footer(_s)
# House slides draw ALL their own furniture — hide the master's shapes (incl. its
# footer text box) on the BLANK layout so nothing phantom leaks into their text layer.
BLANK._element.set("showMasterSp", "0")
import math

SLIDE_MAP={"labs":{}, "domains":{}, "sections":{}, "chapters":{}}

OLD_IDS=list(prs.slides._sldIdLst)   # the legacy slides, in their original order
ORDER=[]                             # final slide sequence (sldId elements)

def slide():
    s=prs.slides.add_slide(BLANK)
    ORDER.append(prs.slides._sldIdLst[-1])
    return s
def rect(s,x,y,w,h,color,line=None):
    sp=s.shapes.add_shape(1,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False; return sp
def oval(s,x,y,w,h,color):
    sp=s.shapes.add_shape(9,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    sp.line.fill.background(); sp.shadow.inherit=False; return sp
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=4):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space)
        for t,sz,col,bold in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=bold
            r.font.color.rgb=col; r.font.name="Arial"
    return tb
def bullets(s,x,y,w,h,items,size=18,color=INK,gap=10,mcolor=BLUE):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap)
        lvl=it[1] if isinstance(it,tuple) else 0
        text=it[0] if isinstance(it,tuple) else it
        r=p.add_run(); r.text=("•  " if lvl==0 else "–  ")+text
        r.font.size=Pt(size if lvl==0 else size-2); r.font.color.rgb=color if lvl==0 else GREY
        r.font.name="Arial"; r.font.bold=(lvl==0 and isinstance(it,tuple) and len(it)>2 and it[2])
    return tb

PAGE={"n":0}
def footer(s):
    PAGE["n"]+=1
    txt(s,Inches(0.4),Inches(7.05),Inches(7.5),Inches(0.35),
        [[(f"{COURSE['short']}  ·  {COURSE['code']}",9,GREY,False)]])
    txt(s,Inches(5.0),Inches(7.05),Inches(3.3),Inches(0.35),
        [[("© 2026 Tertiary Infotech Academy Pte Ltd",9,GREY,False)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(12.4),Inches(7.05),Inches(0.6),Inches(0.35),
        [[(str(PAGE["n"]),9,GREY,False)]],align=PP_ALIGN.RIGHT)
    return PAGE["n"]
_CH_BY_NUM={c[0]:c for c in REFERENCE_CHAPTERS}
def insert_chapters(dnum):
    """Splice this domain's legacy chapters (all slides + graphics, minus the
    dropped admin/lab-pointer slides) into the deck at the current position."""
    for ch in DOMAIN_CHAPTERS[dnum]:
        _,a,b,_,title=_CH_BY_NUM[ch]
        SLIDE_MAP["chapters"][str(ch)]=PAGE["n"]+1
        for i in range(a,b+1):
            if i in REFERENCE_DROP: continue
            ORDER.append(OLD_IDS[i-1]); PAGE["n"]+=1

def head(s,title,kicker=None,kcolor=BLUE):
    rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),Inches(1.55),kcolor)
    if kicker: txt(s,Inches(0.85),Inches(0.5),Inches(11.6),Inches(0.4),[[(kicker,14,kcolor,True)]])
    txt(s,Inches(0.85),Inches(0.9),Inches(11.9),Inches(0.9),[[(title,29,INK,True)]])
    rect(s,Inches(0.85),Inches(1.7),Inches(11.63),Inches(0.02),LINE)
    return s
def _logo(name):
    p=os.path.join(ASSETS,name); return p if os.path.exists(p) else None

# ---------------- slide templates ----------------
def cover():
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),RED); rect(s,0,Inches(7.28),SW,Inches(0.22),INK)
    org=_logo("tertiary-infotech-logo.png")
    if org: s.shapes.add_picture(org,Inches(0.85),Inches(0.7),height=Inches(1.05))
    badge=_logo("comptia-linux-logo.png")
    if badge: s.shapes.add_picture(badge,Inches(10.7),Inches(0.55),height=Inches(1.5))
    txt(s,Inches(0.9),Inches(2.3),Inches(12),Inches(0.6),[[("COURSE SLIDES  ·  WSQ",16,RED,True)]])
    txt(s,Inches(0.9),Inches(2.85),Inches(11.5),Inches(1.9),[[(COURSE["title"],38,INK,True)]])
    rect(s,Inches(0.92),Inches(4.85),Inches(2.4),Inches(0.06),RED)
    txt(s,Inches(0.9),Inches(5.15),Inches(12),Inches(1.4),
        [[(f"WSQ Course Code: {COURSE['code']}  ·  CompTIA Linux+ {COURSE['exam']}",15,GREY,False)],
         [(f"Conducted by {COURSE['org']}  ·  UEN {COURSE['uen']}",14,GREY,False)]],space=6)
    txt(s,Inches(0.9),Inches(6.5),Inches(12),Inches(0.4),[[(f"Trainer: {COURSE['trainer']}  ·  Version {COURSE['version']}  ·  {COURSE['version_date']}",12,GREY,False)]])
    txt(s,Inches(0.9),Inches(6.85),Inches(12),Inches(0.34),[[("© 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved.  ·  www.tertiarycourses.com.sg",10,GREY,False)]])
    PAGE["n"]=1   # the cover is physical slide 1, so subsequent footers/slide_map = actual slide index

def section(kicker,title,n,sub="",record=None):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,RED)
    rect(s,Inches(0.85),Inches(2.5),Inches(0.14),Inches(2.0),TEAL)
    txt(s,Inches(1.25),Inches(2.55),Inches(11),Inches(0.6),[[(kicker,18,RED,True)]])
    txt(s,Inches(1.25),Inches(3.0),Inches(11.4),Inches(1.6),[[(title,40,INK,True)]])
    if sub: txt(s,Inches(1.27),Inches(4.55),Inches(11),Inches(0.8),[[(sub,16,GREY,False)]])
    txt(s,Inches(10.0),Inches(0.7),Inches(2.8),Inches(1.6),[[(n,72,RGBColor(0xE2,0xE8,0xF0),True)]],align=PP_ALIGN.RIGHT)
    p=footer(s)
    if record is not None: SLIDE_MAP["domains"][record]=p
def content(title,items,kicker=None,size=20,kcolor=BLUE):
    s=head(slide(),title,kicker,kcolor); bullets(s,Inches(0.85),Inches(1.95),Inches(11.6),Inches(4.9),items,size=size); footer(s); return s
def two_col(title,left,right,kicker=None,lhead="",rhead=""):
    s=head(slide(),title,kicker)
    rect(s,Inches(0.85),Inches(1.95),Inches(5.7),Inches(4.7),LIGHT); rect(s,Inches(6.95),Inches(1.95),Inches(5.55),Inches(4.7),LIGHT)
    if lhead: txt(s,Inches(1.1),Inches(2.15),Inches(5.2),Inches(0.4),[[(lhead,16,BLUE,True)]])
    if rhead: txt(s,Inches(7.2),Inches(2.15),Inches(5.0),Inches(0.4),[[(rhead,16,TEAL,True)]])
    bullets(s,Inches(1.1),Inches(2.7),Inches(5.2),Inches(3.8),left,size=16)
    bullets(s,Inches(7.2),Inches(2.7),Inches(5.05),Inches(3.8),right,size=16,mcolor=TEAL); footer(s); return s
def cards3(title,cards,kicker):
    s=head(slide(),title,kicker); xs=[Inches(0.85),Inches(5.0),Inches(9.15)]
    for i,c in enumerate(cards[:3]):
        x=xs[i]; col=c[0]
        rect(s,x,Inches(1.95),Inches(3.65),Inches(4.7),LIGHT); rect(s,x,Inches(1.95),Inches(3.65),Inches(0.12),col)
        txt(s,x+Inches(0.25),Inches(2.2),Inches(3.2),Inches(0.6),[[(c[1],19,col,True)]])
        bullets(s,x+Inches(0.25),Inches(2.95),Inches(3.2),Inches(3.4),c[2],size=14,mcolor=col,gap=9)
    footer(s); return s
def big_statement(line1,line2,kicker,color=BLUE):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,color)
    txt(s,Inches(1.1),Inches(2.2),Inches(11),Inches(0.5),[[(kicker,16,color,True)]])
    txt(s,Inches(1.1),Inches(2.8),Inches(11.3),Inches(2.4),[[(line1,38,INK,True)]])
    if line2: txt(s,Inches(1.12),Inches(4.9),Inches(11),Inches(1.2),[[(line2,20,GREY,False)]])
    footer(s); return s
def tile_grid(title,items,kicker=None,cols=2,size=15,icons=None,accent=BLUE):
    s=head(slide(),title,kicker,kcolor=accent)
    n=len(items); rows=math.ceil(n/cols)
    X0=Inches(0.85); Y0=Inches(1.95); TOTW=Inches(11.63); AREAH=Inches(4.78)
    gx=Inches(0.3); gy=Inches(0.26)
    cw=int((TOTW-gx*(cols-1))/cols); ch=int((AREAH-gy*(rows-1))/rows)
    bd=Inches(0.6)
    for i,it in enumerate(items):
        r=i//cols; c=i%cols
        x=int(X0+(cw+gx)*c); y=int(Y0+(ch+gy)*r); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,Inches(0.1),ch,col)
        oval(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,col)
        ic=icons[i] if icons else str(i+1)
        txt(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,[[(ic,19,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        tx=x+Inches(1.08); tw=cw-Inches(1.32)
        if isinstance(it,tuple):
            txt(s,tx,int(y+Inches(0.14)),tw,int(ch-Inches(0.2)),
                [[(it[0],size+2,INK,True)],[(it[1],size-2,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE,space=3)
        else:
            txt(s,tx,int(y+Inches(0.1)),tw,int(ch-Inches(0.16)),[[(it,size,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def flow_h(title,steps,kicker=None,color=RED):
    s=head(slide(),title,kicker,kcolor=color)
    n=len(steps); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.34)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(2.55); ch=Inches(3.15); bd=Inches(0.82)
    for i,st in enumerate(steps):
        x=int(X0+(cw+gap)*i)
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.1),color)
        oval(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,color)
        txt(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,[[(str(i+1),30,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+Inches(0.16),int(y+Inches(1.55)),cw-Inches(0.32),int(ch-Inches(1.7)),[[(st,14,INK,False)]],align=PP_ALIGN.CENTER)
        if i<n-1:
            txt(s,int(x+cw-Inches(0.04)),int(y+ch/2-Inches(0.3)),int(gap+Inches(0.08)),Inches(0.6),
                [[("▶",15,color,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def trainer_slide(kicker,name,role,rows,initials,accent=BLUE):
    s=head(slide(),"About the Trainer",kicker,kcolor=accent)
    lx=Inches(0.85); lw=Inches(3.65)
    rect(s,lx,Inches(1.95),lw,Inches(4.7),LIGHT); rect(s,lx,Inches(1.95),lw,Inches(0.12),accent)
    bd=Inches(1.7); ax=int(lx+(lw-bd)/2)
    oval(s,ax,Inches(2.5),bd,bd,accent)
    txt(s,ax,Inches(2.5),bd,bd,[[(initials,44,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,lx+Inches(0.15),Inches(4.55),lw-Inches(0.3),Inches(0.6),[[(name,21,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,lx+Inches(0.15),Inches(5.2),lw-Inches(0.3),Inches(1.2),[[(role,13,GREY,False)]],align=PP_ALIGN.CENTER)
    rx=Inches(4.9); rw=Inches(7.6); ry=Inches(1.95); rh=Inches(4.7)
    n=len(rows); gy=Inches(0.2); th=int((rh-gy*(n-1))/n)
    for i,(label,val) in enumerate(rows):
        y=int(ry+(th+gy)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,rx,y,rw,th,LIGHT); rect(s,rx,y,Inches(0.1),th,col)
        vruns=[(val,14,INK,False)] if val else [("____________________________________________",13,LINE,False)]
        txt(s,rx+Inches(0.32),y,rw-Inches(0.6),th,
            [[(label.upper(),11,col,True)],vruns],anchor=MSO_ANCHOR.MIDDLE,space=3)
    footer(s); return s
def activity_overview(tag,title,desc,build,services,kicker,record=None):
    s=head(slide(),title,kicker,kcolor=TEAL)
    rect(s,Inches(0.85),Inches(1.85),Inches(1.7),Inches(0.5),TEAL)
    txt(s,Inches(0.85),Inches(1.9),Inches(1.7),Inches(0.4),[[(tag,16,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.85),Inches(2.55),Inches(11.7),Inches(1.6),[[(desc,20,INK,False)]])
    rect(s,Inches(0.85),Inches(4.3),Inches(11.7),Inches(2.0),LIGHT)
    txt(s,Inches(1.1),Inches(4.5),Inches(11),Inches(0.4),[[("You'll build",14,BLUE,True)]])
    txt(s,Inches(1.1),Inches(4.9),Inches(11),Inches(0.6),[[(build,17,INK,True)]])
    txt(s,Inches(1.1),Inches(5.6),Inches(11.2),Inches(0.6),[[("Key commands:  ",13,GREY,True),(services,13,GREY,False)]])
    p=footer(s)
    if record is not None: SLIDE_MAP["labs"][record]=p
    return s
def step_slide(kicker,act_title,n,total,text,cmd=""):
    s=head(slide(),act_title,kicker,TEAL)
    oval(s,Inches(0.85),Inches(2.5),Inches(1.4),Inches(1.4),TEAL)
    txt(s,Inches(0.85),Inches(2.74),Inches(1.4),Inches(0.9),[[(str(n),38,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.95),Inches(1.95),Inches(11),Inches(0.4),[[(f"STEP {n} OF {total}",13,GREY,True)]])
    txt(s,Inches(2.55),Inches(2.4),Inches(10.1),Inches(1.3),[[(text,22,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    if cmd:
        rect(s,Inches(2.55),Inches(4.15),Inches(10.1),Inches(0.95),RGBColor(0x0B,0x12,0x20))
        txt(s,Inches(2.8),Inches(4.28),Inches(9.7),Inches(0.7),[[("$ "+cmd,13,RGBColor(0x9C,0xDC,0xFE),False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def test_slide(act_title,text,kicker):
    s=head(slide(),act_title,kicker,TEAL)
    rect(s,Inches(0.85),Inches(2.3),Inches(11.7),Inches(2.6),RGBColor(0xE8,0xF7,0xEE))
    txt(s,Inches(1.2),Inches(2.6),Inches(11),Inches(0.5),[[("✅  Test it",20,RGBColor(0x12,0x7A,0x3E),True)]])
    txt(s,Inches(1.2),Inches(3.3),Inches(11),Inches(1.4),[[(text,18,INK,False)]]); footer(s); return s
def link_run(p, text, url, size=15, color=None, bold=True, font="Arial"):
    """Add a run that is a REAL clickable hyperlink."""
    r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.bold=bold; r.font.name=font
    r.font.color.rgb=color or BLUE
    r.hyperlink.address=url
    return r
def labs_access(repo, kc):
    """Slide telling learners how to access all the labs from the course GitHub repo,
    with the repo URL as a real clickable hyperlink."""
    s=head(slide(),"Access the Hands-On Labs",kicker="ACCESS THE LABS · GITHUB",kcolor=RED)
    # Repo panel with the clickable link
    rect(s,Inches(0.85),Inches(1.95),Inches(11.63),Inches(1.7),LIGHT)
    rect(s,Inches(0.85),Inches(1.95),Inches(0.12),Inches(1.7),RED)
    txt(s,Inches(1.15),Inches(2.15),Inches(11),Inches(0.5),
        [[("All 30 hands-on labs are on the course GitHub repository — access them by cloning or downloading:",16,INK,True)]])
    lb=s.shapes.add_textbox(Inches(1.15),Inches(2.75),Inches(11.1),Inches(0.7)); lf=lb.text_frame; lf.word_wrap=True
    lp=lf.paragraphs[0]
    link_run(lp, repo, repo, size=17, color=BLUE, font="Consolas")   # clickable repo link
    # Two how-to cards
    cardw=Inches(5.66); y=Inches(4.0); ch=Inches(2.15)
    for i,(col,ttl,lines) in enumerate([
        (BLUE,"Option A — Clone with Git",
         [f"git clone {repo}.git","cd "+repo.rsplit('/',1)[-1]+"/labs","Each lab has its own folder — open its README.md"]),
        (TEAL,"Option B — Download the ZIP",
         ["On GitHub click the green Code button","Choose Download ZIP and unzip it","Open the labs/ folder — one folder per lab"])]):
        x=Inches(0.85)+i*(cardw+Inches(0.31))
        rect(s,x,y,cardw,ch,LIGHT); rect(s,x,y,cardw,Inches(0.1),col)
        txt(s,x+Inches(0.25),y+Inches(0.2),cardw-Inches(0.5),Inches(0.4),[[(ttl,15,col,True)]])
        bullets(s,x+Inches(0.25),y+Inches(0.7),cardw-Inches(0.5),Inches(1.3),lines,size=12,mcolor=col,gap=6)
    txt(s,Inches(0.85),Inches(6.35),Inches(11.6),Inches(0.5),
        [[("Every lab runs free in the browser on Killercoda:  ",12,GREY,True),(kc,12,GREY,False)]])
    footer(s)
def brk(kind,dur,color=AMBER):
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),color); rect(s,0,Inches(7.28),SW,Inches(0.22),color)
    rect(s,Inches(5.4),Inches(2.35),Inches(2.53),Inches(0.1),color)
    txt(s,0,Inches(2.75),SW,Inches(1.2),[[(kind,48,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,0,Inches(4.05),SW,Inches(0.8),[[(dur,22,color,True)]],align=PP_ALIGN.CENTER); PAGE["n"]+=1

def first_cmd(code):
    """Pick a representative single command line from a step's code block."""
    for ln in code.split("\n"):
        s=ln.strip()
        if s and not s.startswith("#") and not s.startswith("cat <<") and "EOF" not in s:
            return s[:88]
    return code.split("\n")[0][:88]

ATTEND=["It is mandatory to take the AM, PM and Assessment digital attendance for WSQ-funded courses.",
 "The trainer/administrator displays the digital attendance QR code from the SSG portal.",
 "Scan the QR code with your mobile phone camera and submit your attendance.",
 "A minimum of 75% attendance is required to be eligible for assessment and funding."]
AFLOW=["TRAQOM digital attendance — scan the SSG QR on the LMS",
 "Assessment digital attendance",
 "Sit WA (SAQ) then PP — open book",
 "Submit your answers on the LMS",
 "Sign the Assessment Summary Record"]

# ============================================================ BUILD
cover()

# ---------------- ADMIN ----------------
section("COURSE ADMINISTRATION","Welcome & Housekeeping","")
tile_grid("Digital Attendance (Mandatory)",ATTEND,kicker="TRAQOM · SSG DIGITAL ATTENDANCE",cols=1,size=13,accent=RED)
trainer_slide("YOUR TRAINER · GENERAL","Your Trainer","General Trainer template —\nto be completed by the trainer",
 [("Name",""),("Title / Designation",""),("Qualifications",""),
  ("Areas of expertise",""),("Training & industry experience",""),("Contact","")],
 initials="?",accent=GREY)
trainer_slide("YOUR TRAINER",COURSE["trainer"],"Principal Trainer\nTertiary Infotech Academy Pte. Ltd.",
 [("Role","Principal Trainer, Tertiary Infotech Academy Pte. Ltd."),
  ("Certification","Linux / cloud / DevOps & cybersecurity — hands-on Linux systems administration."),
  ("Delivers","WSQ courses on Linux, cloud administration, DevOps and software engineering."),
  ("Founder","Founder and lead instructor at Tertiary Infotech / Tertiary Courses.")],
 initials="AA",accent=RED)
content("Let's Know Each Other",[
 "Your name and organisation / role.",
 "Your current experience with Linux (beginner → advanced).",
 "Which of the five exam domains matters most to your day job, and your exam timeline."],kicker="ICE-BREAKER")
tile_grid("Ground Rules",[
 "Set your mobile phone to silent mode.","Participate actively — no question is too small.",
 "Mutual respect: agree to disagree.","Do every lab yourself on Killercoda.",
 "Be punctual; return from breaks on time.","75% attendance is required."],
 kicker="HOUSEKEEPING",cols=2,size=15,accent=RED)
# Download Course Material — VISUAL slide: portal screenshot + numbered steps
def download_material():
    s=head(slide(),"Download Course Material",kicker="LMS / TMS · COURSE PORTAL",kcolor=BLUE)
    shot=_logo("lms-tms-portal.png")
    if shot:
        pic=s.shapes.add_picture(shot,Inches(6.6),Inches(2.0),width=Inches(6.2))
        pic.line.color.rgb=LINE; pic.line.width=_EmuPt(1)
    steps=[("1","Go to the portal","lms-tms.tertiaryinfotech.com"),
           ("2","Log in","Use your registered email — Send OTP (or password)."),
           ("3","Open your course","CompTIA Certified Linux+ Training (TGS-2024048316)."),
           ("4","Download the materials","Slides, Learner Guide and Lesson Plan — allowed in the open-book assessment.")]
    y=Inches(2.0); hstep=Inches(1.12); bd=Inches(0.5)
    for i,(nn,t1,t2) in enumerate(steps):
        yy=int(y+hstep*i)
        rect(s,Inches(0.85),yy,Inches(5.5),int(hstep-Inches(0.14)),LIGHT)
        rect(s,Inches(0.85),yy,Inches(0.1),int(hstep-Inches(0.14)),PALETTE[i%len(PALETTE)])
        oval(s,Inches(1.05),int(yy+(hstep-Inches(0.14))/2-bd/2),bd,bd,PALETTE[i%len(PALETTE)])
        txt(s,Inches(1.05),int(yy+(hstep-Inches(0.14))/2-bd/2),bd,bd,[[(nn,15,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,Inches(1.75),yy,Inches(4.5),int(hstep-Inches(0.14)),
            [[(t1,14,INK,True)],[(t2,11,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE,space=2)
    txt(s,Inches(6.6),Inches(6.35),Inches(6.2),Inches(0.4),
        [[(f"Portal: {COURSE['lms']}",12,GREY,False)]],align=PP_ALIGN.CENTER)
    footer(s)
download_material()
two_col("Lesson Plan — 2 Days, 8 hours/day",[
 ("Day 1 — System & User Management",0),
 ("Domain 1: System Management, 23% (Labs 1–7)",1),
 ("Domain 2: Services & User Management, 20% (Labs 8–13)",1),
 ("9:30am–6:30pm · 1-hour lunch · tea within",1)],
 [("Day 2 — Security, Automation & Troubleshooting",0),
 ("Domain 3: Security, 18% (Labs 14–19)",1),
 ("Domain 4: Automation & Scripting, 17% (Labs 20–24)",1),
 ("Domain 5: Troubleshooting, 22% (Labs 25–30)",1),
 ("Final Assessment (WA + PP) from 4:30pm",1)],
 kicker="SCHEDULE",lhead="Day 1",rhead="Day 2")
tile_grid("Learning Outcomes",[DOMAIN_OUTCOMES[d["num"]] for d in DOMAINS],
 kicker="WHAT YOU'LL ACHIEVE",cols=1,size=14,accent=BLUE)
# The XK0-006 exam blueprint, shown at the START of the deck — the whole course
# is sequenced on these five domains.
tile_grid("CompTIA Linux+ XK0-006 — Exam Domains",
 [(f"Domain {d['num']} — {d['title']}",
   f"{d['weight']}% of the exam  ·  Objectives {d['objs'][0][0]}–{d['objs'][-1][0]}  ·  Labs {d['labs'][0]}–{d['labs'][-1]}")
  for d in DOMAINS],
 kicker="EXAM BLUEPRINT · XK0-006 V8",cols=1,size=14,accent=RED)
content("How This Deck Is Organised",[
 "The course follows the five XK0-006 exam domains in blueprint order.",
 "Each domain: exam objectives → the full teaching chapters (theory + examples) → hands-on labs → recap.",
 "All teaching content and graphics from the course study chapters are included, re-ordered onto the exam domains.",
 "Every lab maps to a specific exam objective — shown on each lab's overview slide."],
 kicker="EXAM-ALIGNED STRUCTURE")
tile_grid("Briefing for Assessment",[
 "Place phones and other materials under the table or on the floor.",
 "No photos or recording of assessment scripts.","No discussion during the assessment.",
 "Use a black/blue pen for hard-copy assessments.","No liquid paper / correction tape.",
 "Scripts are collected when time is up."],kicker="BEFORE THE ASSESSMENT",cols=2,size=12,accent=RED)
tile_grid("Assessment",[
 ("Written Assessment (WA)","Short-Answer Questions (SAQ) · 1 hour · open book — aligned to the slides."),
 ("Practical Performance (PP)","Hands-on Linux tasks · 1 hour · open book — aligned to the labs."),
 ("Open Book","Slides, Learner Guide and approved materials only."),
 ("Eligibility","Minimum 75% attendance for assessment and funding · an appeal process is available.")],
 kicker="FINAL ASSESSMENT",cols=1,size=13,accent=RED)
flow_h("Assessment Flow",AFLOW,kicker="ON ASSESSMENT DAY")

# ---------------- CORE LINUX CONCEPTS ----------------
section("CORE CONCEPTS","Linux Fundamentals for Administrators","")
tile_grid("What is Linux?",[
 ("Open-source OS","A free, open-source operating system kernel plus GNU userland, powering most servers and clouds."),
 ("Distributions","Debian/Ubuntu (apt/dpkg) and RHEL/Rocky/Fedora (dnf/rpm) package the kernel differently."),
 ("Many ways to manage","The shell (bash), config files under /etc, systemd, and automation with Ansible/Bash/Python."),
 ("The Linux+ role","Configure, manage, secure, automate and troubleshoot Linux servers.")],
 kicker="OVERVIEW",cols=2,size=15)
cards3("The Linux Landscape",[
 (BLUE,"Kernel & FHS",["Kernel + initramfs boot userspace","FHS: /etc /var /usr /home …","Everything is a file"]),
 (TEAL,"Distributions",["Debian/Ubuntu — apt/dpkg","RHEL/Rocky/Fedora — dnf/rpm","Same kernel, different tooling"]),
 (VIOLET,"Architectures & licensing",["x86_64, ARM64 (AArch64), RISC-V","Free / open-source vs proprietary","Copyleft (GPL) licensing"])],kicker="THE ECOSYSTEM")
big_statement("On Linux, everything is a file.","Files, directories, devices, sockets and processes are all reachable through the filesystem — master the shell and you master the system.","WHY IT MATTERS",color=BLUE)
two_col("Ways to Work on Linux",[
 ("The shell (bash)",0),("Interactive commands + pipes",1),("Fast and scriptable",1),
 ("Text editors",0),("vi/vim and nano",1),("Edit config files in place",1)],
 [("Scripting & automation",0),("Bash and Python scripts",1),("Ansible for Infrastructure as Code",1),
 ("Containers & services",0),("systemd manages services",1),("Docker/Podman run apps",1)],
 kicker="TOOLING",lhead="Interactive",rhead="Automation")
content("Your Workbench — Killercoda",[
 f"Every lab runs free in the browser on the Killercoda Ubuntu Playground: {COURSE['killercoda']}",
 "A disposable Ubuntu VM — install packages, break things and reset with no risk.",
 "Type the commands yourself; muscle memory is what the exam tests.",
 "Reset the playground between labs that change kernel, firewall or systemd state."],kicker="HANDS-ON")
labs_access(COURSE["repo"], COURSE["killercoda"])

# ---------------- DOMAINS + LABS ----------------
CARD_COLORS=[BLUE,TEAL,VIOLET]
for d in DOMAINS:
    code=f"{d['num']:02d}"
    sub=f"Objectives {d['objs'][0][0]}–{d['objs'][-1][0]}  ·  {d['weight']}% of the exam  ·  Labs {d['labs'][0]}–{d['labs'][-1]}"
    section(f"DOMAIN {code}", d["title"], code, sub, record=d["num"])
    tile_grid(f"Key Concepts — {d['title']}", DOMAIN_CONCEPTS[d["num"]],
              kicker=f"EXAM WEIGHTING {d['weight']}%", cols=1, size=14, accent=CARD_COLORS[0])
    tile_grid(f"Exam Objectives — {d['title']}",
              [(f"Objective {o[0]}", o[1]) for o in d["objs"]],
              kicker=f"DOMAIN {code} · BLUEPRINT", cols=1, size=13, accent=RED)
    # Full legacy teaching chapters for this domain (all content + graphics,
    # re-ordered from the v3 deck onto this exam domain).
    insert_chapters(d["num"])
    acts=domain_labs(d["num"])
    third=(len(acts)+2)//3
    groups=[acts[i:i+third] for i in range(0,len(acts),third)][:3]
    while len(groups)<3: groups.append([])
    cards=[]
    for gi,g in enumerate(groups):
        if not g: label="—"
        elif len(g)==1: label=f"Lab {g[0]['num']}"
        else: label=f"Labs {g[0]['num']}–{g[-1]['num']}"
        cards.append((CARD_COLORS[gi], label, [a["title"] for a in g] if g else ["—"]))
    cards3(f"Hands-On Labs — {d['title']}", cards, kicker="WHAT YOU'LL DO")
    for a in acts:
        services=", ".join(a["concepts"][:6])
        activity_overview(f"LAB {a['num']}", a["title"], a["goal"], a["build"], services,
                          kicker=f"DOMAIN {code} · OBJ {a['objective']} · HANDS-ON", record=a["num"])
        steps=a["steps"]; total=len(steps)
        for i,st in enumerate(steps,1):
            step_slide(f"LAB {a['num']} · {a['title'][:32]}", a["title"], i, total,
                       st["title"], first_cmd(st["code"]))
        test_slide(a["title"], a["test"], kicker=f"LAB {a['num']} · VERIFY")
    content(f"Recap — {d['title']}",
            [f"Lab {a['num']}: {a['title']}" for a in acts],
            kicker="DOMAIN RECAP", size=16)
    # break dividers aligned to the 2-day Lesson Plan (uses the reference `brk` divider)
    BREAKS = {1: ("Tea Break", "15 minutes", AMBER),
              2: ("End of Day 1", "See you tomorrow — 9:00 AM", TEAL),
              3: ("Lunch Break", "1 hour", AMBER),
              4: ("Tea Break", "15 minutes", AMBER)}
    if d["num"] in BREAKS:
        brk(*BREAKS[d["num"]])

# ---------------- CLOSE ----------------
section("WRAP-UP","Course Summary & Next Steps","")
tile_grid("What You Achieved",[DOMAIN_OUTCOMES[d["num"]] for d in DOMAINS],
 kicker="LEARNING OUTCOMES",cols=1,size=14)
content("Preparing for the Linux+ XK0-006 Exam",[
 "Redo every lab on Killercoda until the commands are automatic.",
 "Review the 'Test it' takeaway in each lab and the Learner Guide.",
 "Know which tool solves which problem — the exam is scenario-based (performance-based items).",
 "The exam has up to 90 questions (multiple-choice + performance-based) in 90 minutes.",
 "Book the exam through Pearson VUE from your CompTIA account."],kicker="NEXT STEPS")
# CompTIA Linux+ exam registration — shown at the END of the deck (house request)
flow_h("CompTIA Linux+ Exam Registration",[
 "Create your CompTIA account at comptia.org",
 "Buy the XK0-006 exam voucher (CompTIA Store / Tertiary Infotech)",
 "Schedule at a Pearson VUE test centre or online proctored",
 "Sit the exam — up to 90 questions in 90 minutes, pass 720/900",
 "Claim your digital badge and certificate"],kicker="GET CERTIFIED · XK0-006")
tile_grid("Exam Registration — Links & Details",[
 ("Exam voucher",f"CompTIA Store: {COURSE['exam_voucher']}  (or through Tertiary Infotech)"),
 ("Schedule the exam",f"Pearson VUE: {COURSE['exam_pearson']} — test centre or online proctored"),
 ("Exam format","XK0-006 · up to 90 questions · 90 minutes · passing score 720 (scale 100–900)."),
 ("Question types","Multiple-choice + performance-based items (live Linux tasks)."),
 ("On exam day","Two forms of ID at a test centre · check the online-testing system requirements for remote.")],
 kicker="REGISTER FOR THE EXAM",cols=1,size=12,accent=RED)
tile_grid("Practice Exam",[
 ("Test yourself","Sharpen your exam readiness with the Tertiary Infotech CompTIA Linux+ practice exam."),
 ("Practice exam link",COURSE['practice_exam']),
 ("Timed conditions","Attempt it under timed conditions and review every explanation."),
 ("Close the gaps","Revisit any lab whose domain you miss, then re-take the practice exam.")],
 kicker="TEST YOURSELF · EXAMS.TERTIARYINFOTECH.COM",cols=1,size=13,accent=TEAL)
# Assessment admin repeated at the END (house rule)
tile_grid("Assessment",[
 ("Two instruments","Written Assessment (SAQ) — 1 hour · Practical Performance (PP) — 1 hour."),
 ("Open book","Slides, Learner Guide and approved materials only."),
 ("Digital attendance","Remember the Assessment digital attendance (TRAQOM · SSG)."),
 ("Submission",f"Submit your completed answers on the LMS at {COURSE['lms']}")],
 kicker="WRAP-UP",cols=1,size=13,accent=RED)
flow_h("Assessment Flow",AFLOW,kicker="ON ASSESSMENT DAY")
tile_grid("Digital Attendance (Mandatory)",ATTEND,kicker="TRAQOM · SSG DIGITAL ATTENDANCE",cols=1,size=13,accent=RED)
big_statement("Thank You!","You are now ready to administer Linux — and to sit the CompTIA Linux+ XK0-006 exam.","ALL THE BEST",color=TEAL)

SLIDE_MAP["total"]=PAGE["n"]

# ---------------- final re-sequencing ----------------
# Drop the legacy slides that were not carried over (superseded admin block,
# part dividers, retired lab pointers, old close), then move every slide into
# its final exam-domain position. Re-appending an element that already has a
# parent MOVES it, so iterating ORDER yields exactly the ORDER sequence.
sldIdLst=prs.slides._sldIdLst
_keep={id(el) for el in ORDER}
for el in OLD_IDS:
    if id(el) not in _keep:
        prs.part.drop_rel(el.rId)
        sldIdLst.remove(el)
for el in ORDER:
    sldIdLst.append(el)
assert len(sldIdLst)==len(ORDER)==PAGE["n"], \
    f"slide count mismatch: sldIdLst={len(sldIdLst)} ORDER={len(ORDER)} PAGE={PAGE['n']}"

OUT=os.path.join(COURSEWARE,f"PPT-CompTIA-Linux-Plus-XK0-006-{COURSE['version']}.pptx")
prs.save(OUT)
with open(os.path.join(COURSEWARE,"slide_map.json"),"w") as fh:
    json.dump(SLIDE_MAP,fh,indent=2)
print(f"Saved {OUT}  ({PAGE['n']} slides)")
print("Wrote slide_map.json")
